import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import os

from fpdf import FPDF
from PIL import Image
from moviepy import AudioFileClip, VideoFileClip
from docx import Document
from PyPDF2 import PdfReader
from pdf2image import convert_from_path
import pandas as pd
from pptx import Presentation
from ebooklib import epub

# Supported formats
conversions = {
    ".docx": [".pdf", ".txt"],
    ".pdf": [".docx", ".jpg", ".png"],
    ".jpg": [".pdf"],
    ".png": [".pdf"],
    ".mp4": [".mp3"],
    ".wav": [".mp3"],
    ".mp3": [".wav"],
    ".txt": [".pdf"],
    ".xlsx": [".csv"],
    ".csv": [".xlsx"],
    ".pptx": [".pdf"],
    ".epub": [".pdf"],
}


def convert_file(filepath, target_ext):
    file_ext = os.path.splitext(filepath)[1].lower()
    output_path = os.path.splitext(filepath)[0] + "_converted" + target_ext

    try:
        # DOCX → PDF
        if file_ext == ".docx" and target_ext == ".pdf":
            doc = Document(filepath)
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            for para in doc.paragraphs:
                pdf.multi_cell(0, 10, para.text)
            pdf.output(output_path)

        # DOCX → TXT
        elif file_ext == ".docx" and target_ext == ".txt":
            doc = Document(filepath)
            with open(output_path, "w", encoding="utf-8") as f:
                for para in doc.paragraphs:
                    f.write(para.text + "\n")

        # PDF → DOCX
        elif file_ext == ".pdf" and target_ext == ".docx":
            reader = PdfReader(filepath)
            doc = Document()
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    doc.add_paragraph(text)
            doc.save(output_path)

        # PDF → JPG/PNG (one image per page)
        elif file_ext == ".pdf" and target_ext in [".jpg", ".png"]:
            images = convert_from_path(filepath)
            for i, img in enumerate(images):
                img_path = output_path.replace(target_ext, f"_{i+1}{target_ext}")
                img.save(img_path)

        # Image → PDF
        elif file_ext in [".jpg", ".png"] and target_ext == ".pdf":
            img = Image.open(filepath).convert("RGB")
            img.save(output_path)

        # MP4/WAV → MP3
        elif file_ext in [".mp4", ".wav"] and target_ext == ".mp3":
            clip = (
                AudioFileClip(filepath)
                if file_ext == ".wav"
                else VideoFileClip(filepath).audio
            )
            clip.write_audiofile(output_path)

        # MP3 → WAV
        elif file_ext == ".mp3" and target_ext == ".wav":
            clip = AudioFileClip(filepath)
            clip.write_audiofile(output_path)

        # TXT → PDF
        elif file_ext == ".txt" and target_ext == ".pdf":
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            with open(filepath, "r", encoding="utf-8") as f:
                for line in f:
                    pdf.multi_cell(0, 10, line.strip())
            pdf.output(output_path)

        # XLSX → CSV
        elif file_ext == ".xlsx" and target_ext == ".csv":
            df = pd.read_excel(filepath)
            df.to_csv(output_path, index=False)

        # CSV → XLSX
        elif file_ext == ".csv" and target_ext == ".xlsx":
            df = pd.read_csv(filepath)
            df.to_excel(output_path, index=False)

        # PPTX → PDF (extracts text)
        elif file_ext == ".pptx" and target_ext == ".pdf":
            prs = Presentation(filepath)
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.set_font("Arial", size=12)
            for slide in prs.slides:
                pdf.add_page()
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        pdf.multi_cell(0, 10, shape.text)
            pdf.output(output_path)

        # EPUB → PDF (extracts all document text)
        elif file_ext == ".epub" and target_ext == ".pdf":
            book = epub.read_epub(filepath)
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            for item in book.get_items():
                if item.get_type() == epub.ITEM_DOCUMENT:
                    text = item.get_content().decode("utf-8")
                    pdf.multi_cell(0, 10, text)
            pdf.output(output_path)

        else:
            messagebox.showerror(
                "Error", f"Conversion from {file_ext} to {target_ext} is not supported."
            )
            return

        messagebox.showinfo("Success", f"File converted and saved as:\n{output_path}")

    except Exception as e:
        messagebox.showerror("Conversion Error", str(e))


# GUI
def browse_file():
    filepath = filedialog.askopenfilename()
    if filepath:
        file_entry.delete(0, tk.END)
        file_entry.insert(0, filepath)
        ext = os.path.splitext(filepath)[1].lower()
        format_menu["values"] = conversions.get(ext, [])
        if conversions.get(ext):
            format_menu.current(0)
        else:
            messagebox.showwarning(
                "Unsupported", f"No supported conversions for: {ext}"
            )


def handle_convert():
    filepath = file_entry.get()
    target_ext = format_var.get()
    if not filepath or not os.path.exists(filepath):
        messagebox.showerror("No File", "Please select a valid file.")
        return
    if not target_ext:
        messagebox.showerror("No Format", "Please select a target format.")
        return
    convert_file(filepath, target_ext)


# App Window
app = tk.Tk()
app.title("File Format Converter")
app.geometry("500x250")
app.resizable(False, False)

tk.Label(app, text="Select File:").pack(pady=(20, 5))
file_frame = tk.Frame(app)
file_frame.pack()
file_entry = tk.Entry(file_frame, width=50)
file_entry.pack(side=tk.LEFT, padx=5)
tk.Button(file_frame, text="Browse", command=browse_file).pack(side=tk.LEFT)

tk.Label(app, text="Convert To:").pack(pady=10)
format_var = tk.StringVar()
format_menu = ttk.Combobox(app, textvariable=format_var, state="readonly", width=20)
format_menu.pack()

tk.Button(
    app,
    text="Convert",
    command=handle_convert,
    bg="#4CAF50",
    fg="white",
    padx=10,
    pady=5,
).pack(pady=20)

app.mainloop()
